import json
import os
import time
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import tempfile
import boto3
import pytesseract
from PIL import Image
import io
from typing import Dict, Optional, Any
from db_query import database_query

# S3 bucket fallback for local testing
PRIVATE_BUCKET = os.environ.get("PRIVATE_FILES_BUCKET", "private-bucket-fastapi")

# Tesseract OCR path (for Lambda environment)
pytesseract.pytesseract.tesseract_cmd = "/usr/bin/tesseract"

# === AWS Bedrock Setup ===
bedrock = boto3.client("bedrock-runtime", region_name="us-east-1")
TEXT_MODEL_ARN = "arn:aws:bedrock:us-east-1:127214171089:inference-profile/us.meta.llama4-scout-17b-instruct-v1:0"


def _parse_db_result(result):
    """Helper to parse database_query results"""
    # Check for errors first
    if isinstance(result, dict) and result.get('statusCode') == 500:
        print(f"Database query error: {result.get('error')}")
        return []
    
    # Handle different response formats
    if isinstance(result, str):
        try:
            body = json.loads(result)
        except json.JSONDecodeError:
            print(f"Failed to parse string result: {result}")
            return []
    elif isinstance(result, dict) and "body" in result:
        try:
            body = json.loads(result["body"]) if isinstance(result["body"], str) else result["body"]
        except json.JSONDecodeError:
            print(f"Failed to parse body: {result['body']}")
            return []
    else:
        body = result
    
    if not isinstance(body, list):
        print(f"Expected list but got: {type(body)} - {body}")
        return []
    
    return body


def invoke_bedrock_text(system_msg: str, user_content: str, temperature: float = 0.3, max_tokens: int = 2048) -> str:
    """Invoke Bedrock text model and return response."""
    try:
        response = bedrock.converse(
            modelId=TEXT_MODEL_ARN,
            system=[{"text": system_msg.strip()}] if system_msg else [],
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"text": user_content.strip()}
                    ]
                }
            ],
            inferenceConfig={
                "maxTokens": max_tokens,
                "temperature": temperature,
                "topP": 0.9
            }
        )
        
        if response.get("output") and response["output"].get("message"):
            content = response["output"]["message"].get("content", [])
            if content:
                return content[0].get("text", "").strip()
        return ""
    except Exception as e:
        print(f"Bedrock invocation failed: {e}")
        import traceback
        print(traceback.format_exc())
        return ""


def takeaways_json_to_text(takeaways_json: str) -> str:
    """Convert JSON array string to bullet-point text format."""
    try:
        # Clean up the response (remove markdown code blocks if present)
        cleaned = takeaways_json.replace("```json", "").replace("```", "").strip()
        
        # Parse the JSON array string
        points = json.loads(cleaned)
        
        # Join as bullet points, one per line
        if isinstance(points, list):
            return "\n".join([f"- {point}" if not point.startswith("-") else point for point in points])
        else:
            return str(points)
    except json.JSONDecodeError as e:
        print(f"[WARNING] Failed to parse as JSON, returning as-is: {e}")
        # If not JSON, try to extract list items or return as-is
        return takeaways_json
    except Exception as e:
        print(f"[ERROR] Failed to convert takeaways JSON to text: {e}")
        return takeaways_json


def extract_text_from_file(file_path: str) -> str:
    """Extract text from PDF, DOCX, or PPTX files."""
    ext = os.path.splitext(file_path)[1].lower()
    content = ""
    print(f"[DEBUG] Extracting text from file: {file_path} (ext: {ext})")
    
    try:
        if ext == ".pdf":
            doc = None
            try:
                doc = fitz.open(file_path)
                for page in doc:
                    page_text = page.get_text()
                    if page_text.strip():
                        content += page_text + "\n"
                
                # If no text found, try OCR
                if not content.strip():
                    print("[DEBUG] No text found in PDF, attempting OCR on images...")
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        pix = page.get_pixmap()
                        img_bytes = pix.tobytes("png")
                        img = Image.open(io.BytesIO(img_bytes))
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            content += ocr_text + "\n"
                        img.close()  # Close image
            finally:
                # Ensure document is always closed
                if doc is not None:
                    doc.close()
            
        elif ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    content += para.text + "\n"
            
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                content += text + "\n"
        else:
            raise ValueError(f"Unsupported file type: {ext}. Use PDF, DOCX, or PPTX.")
            
    except Exception as e:
        print(f"[ERROR] Failed to extract text: {e}")
        import traceback
        print(traceback.format_exc())
        raise
    
    print(f"[DEBUG] Extracted text length: {len(content)} characters")
    return content.strip()


def get_takeaways_from_text(text: str) -> str:
    """Generate key takeaways from text using Bedrock."""
    print(f"[DEBUG] Generating takeaways from text (length: {len(text)})")
    
    try:
        # Limit text to avoid token limits (keep first 10000 characters)
        text_for_analysis = text[:10000] if len(text) > 10000 else text
        
        system_msg = (
            "You are an expert in document analysis and summarization. "
            "Extract key takeaways from documents and return them as a JSON array of strings."
        )
        
        prompt = f"""
Analyze the following document content and extract 2-5 key takeaways.

Requirements:
- Return ONLY a JSON array of strings, e.g. ["point 1", "point 2", "point 3"]
- Each point should be concise and informative
- Do not include any explanation, markdown formatting, or code blocks
- Focus on the most important insights, facts, or recommendations
- Return pure JSON only, no other text

Document content:
{text_for_analysis}
"""
        
        print(f"[DEBUG] Invoking Bedrock for takeaways generation...")
        response = invoke_bedrock_text(system_msg, prompt, temperature=0.3, max_tokens=1024)
        
        if not response or not response.strip():
            return "[]"
        
        print(f"[DEBUG] Bedrock response received (length: {len(response)})")
        return response.strip()
        
    except Exception as e:
        print(f"[ERROR] Error in get_takeaways_from_text: {e}")
        import traceback
        print(traceback.format_exc())
        raise


def get_takeaways_from_file(file_path: str) -> str:
    """Extract text from file and generate takeaways."""
    print(f"[DEBUG] Getting takeaways from file: {file_path}")
    
    try:
        text = extract_text_from_file(file_path)
        
        if not text or not text.strip():
            raise ValueError("No text could be extracted from the file")
        
        takeaways_json = get_takeaways_from_text(text)
        return takeaways_json
        
    except Exception as e:
        print(f"[ERROR] Error in get_takeaways_from_file: {e}")
        import traceback
        print(traceback.format_exc())
        raise


def update_fact_pack_highlights(rowid: int, key_highlights: str) -> bool:
    """Update key_highlights in fact_pack table."""
    try:
        print(f"[DEBUG] Updating fact_pack table: id={rowid}, key_highlights length={len(key_highlights)}")
        
        query = "UPDATE fact_pack SET key_highlights = %s WHERE id = %s"
        params = [key_highlights, rowid]
        
        result = database_query(query, params)
        
        # Check if update was successful
        if isinstance(result, dict):
            if result.get('statusCode') == 500:
                print(f"[ERROR] Database update failed: {result.get('error')}")
                return False
            # Check if result indicates success (could be rowcount or success message)
            if result.get('rowcount', 0) == 0 and 'error' not in result:
                print(f"[WARNING] No rows updated. Check if id={rowid} exists in fact_pack.")
                # Still return True as the query executed successfully, just no rows matched
        
        print(f"[DEBUG] Database update successful")
        return True
        
    except Exception as e:
        print(f"[ERROR] Error updating fact_pack: {e}")
        import traceback
        print(traceback.format_exc())
        return False


def process_s3_file(bucket: str, key: str, rowid: int) -> Dict[str, Any]:
    """Process file from S3 and update fact_pack table."""
    s3 = boto3.client('s3')
    ext = os.path.splitext(key)[1].lower()
    
    print(f"[DEBUG] Processing S3 file: s3://{bucket}/{key}, rowid={rowid}")
    
    tmp_path = None
    try:
        # Download file to temp - close the file handle immediately after writing
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            s3.download_fileobj(bucket, key, tmp)
            tmp_path = tmp.name
            # File handle is closed here when exiting the 'with' block
        
        print(f"[DEBUG] Downloaded file to: {tmp_path}")
        
        # Extract takeaways
        takeaways_json = get_takeaways_from_file(tmp_path)
        takeaways_text = takeaways_json_to_text(takeaways_json)
        
        # Update database
        success = update_fact_pack_highlights(rowid, takeaways_text)
        
        if not success:
            return {
                "success": False,
                "error": "Failed to update database"
            }
        
        return {
            "success": True,
            "rowid": rowid,
            "takeaways": takeaways_text
        }
        
    finally:
        # Clean up temp file - make deletion non-critical
        if tmp_path and os.path.exists(tmp_path):
            try:
                # On Windows, sometimes we need to wait a bit for file handles to be released
                time.sleep(0.1)
                os.remove(tmp_path)
                print(f"[DEBUG] Removed temporary file: {tmp_path}")
            except (PermissionError, OSError) as e:
                # File deletion is not critical - just log a warning
                print(f"[WARNING] Could not delete temporary file {tmp_path}: {e}")
                # Try to delete on next attempt (not critical for Lambda)


def lambda_handler(event, context=None):
    """
    Lambda handler for fact-pack processing.
    
    Expected event formats:
    
    1. S3 Event (triggered by S3 upload):
    {
        "Records": [{
            "s3": {
                "bucket": {"name": "bucket-name"},
                "object": {"key": "path/to/file.pdf"}
            }
        }]
    }
    Note: rowid should be in S3 object metadata
    
    2. Direct call:
    {
        "action": "fact_pack",
        "bucket": "bucket-name",
        "key": "path/to/file.pdf",
        "rowid": 123
    }
    """
    try:
        print(f"[DEBUG] Received event: {json.dumps(event, indent=2)}")
        
        # Handle S3 event format
        if "Records" in event and len(event["Records"]) > 0:
            record = event["Records"][0]
            bucket = record["s3"]["bucket"]["name"]
            key = record["s3"]["object"]["key"]
            
            # Get object metadata for rowid
            s3 = boto3.client('s3')
            try:
                obj = s3.head_object(Bucket=bucket, Key=key)
                rowid = obj.get('Metadata', {}).get('rowid')
                
                if not rowid:
                    return {
                        "statusCode": 400,
                        "body": json.dumps({
                            "success": False,
                            "error": "rowid not found in S3 object metadata"
                        })
                    }
                
                rowid = int(rowid)
                
            except Exception as e:
                print(f"[ERROR] Failed to get S3 object metadata: {e}")
                return {
                    "statusCode": 400,
                    "body": json.dumps({
                        "success": False,
                        "error": f"Failed to get S3 object metadata: {str(e)}"
                    })
                }
        
        # Handle direct call format
        elif "bucket" in event and "key" in event and "rowid" in event:
            bucket = event["bucket"]
            key = event["key"]
            rowid = int(event["rowid"])
        
        else:
            return {
                "statusCode": 400,
                "body": json.dumps({
                    "success": False,
                    "error": "Invalid event format. Expected S3 event or direct call with bucket, key, and rowid"
                })
            }
        
        # Process the file
        result = process_s3_file(bucket, key, rowid)
        
        if result.get("success"):
            return {
                "statusCode": 200,
                "body": json.dumps(result)
            }
        else:
            return {
                "statusCode": 500,
                "body": json.dumps(result)
            }
            
    except Exception as e:
        print(f"[ERROR] Lambda handler error: {e}")
        import traceback
        traceback.print_exc()
        return {
            "statusCode": 500,
            "body": json.dumps({
                "success": False,
                "error": f"Lambda execution failed: {str(e)}"
            })
        }


if __name__ == "__main__":
    event = {
        "bucket": PRIVATE_BUCKET,
        "key": "factPack/ce1f6bb8-6388-4926-aec8-2cc3519a1591.pdf",
        "rowid": 1
    }
    result = lambda_handler(event)
    print(result)